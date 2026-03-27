#!/home/tomoyuki/.openclaw/workspace/.venv-pptx/bin/python
from __future__ import annotations

import traceback
from pathlib import Path
from typing import Any, Dict, List

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt

FONT_FAMILY = "Noto Sans CJK JP"

# Colors
BG_DARK = RGBColor(0x0D, 0x1B, 0x2A)         # #0D1B2A
BG_DARK_2 = RGBColor(0x15, 0x22, 0x38)       # #152238
ACCENT_MAIN = RGBColor(0x00, 0xD4, 0xFF)     # #00D4FF
ACCENT_SUB = RGBColor(0x39, 0xFF, 0x14)      # #39FF14
TEXT_LIGHT = RGBColor(0xE0, 0xE0, 0xE0)      # #E0E0E0
TEXT_WHITE = RGBColor(0xFF, 0xFF, 0xFF)      # #FFFFFF
CARD_BG = RGBColor(0x1B, 0x28, 0x38)         # #1B2838


def mm(v: float) -> int:
    return int(Inches(v / 25.4))


def i(v: Any) -> int:
    return int(v)


def apply_run_font(run, *, size_pt: float, color: RGBColor, bold: bool = False, font_name: str = FONT_FAMILY) -> None:
    run.font.name = font_name
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = color

    r = run._r
    rPr = r.get_or_add_rPr()
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = OxmlElement("a:ea")
        rPr.append(ea)
    ea.set("typeface", font_name)


def add_textbox(slide, x, y, w, h):
    return slide.shapes.add_textbox(i(x), i(y), i(w), i(h))


def add_rect(slide, x, y, w, h, fill: RGBColor | None = None, line: RGBColor | None = None, rounded: bool = False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, i(x), i(y), i(w), i(h))
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    return shp


def add_base_bg(slide, prs, gradient: bool = False):
    if gradient:
        add_rect(slide, 0, 0, prs.slide_width, mm(105), fill=BG_DARK)
        add_rect(slide, 0, mm(105), prs.slide_width, mm(105), fill=BG_DARK_2)
    else:
        add_rect(slide, 0, 0, prs.slide_width, prs.slide_height, fill=BG_DARK)

    # top line (2mm)
    add_rect(slide, 0, 0, prs.slide_width, mm(2), fill=ACCENT_MAIN, line=ACCENT_MAIN)


def add_slide_number(slide, prs, text: str):
    tb = add_textbox(slide, prs.slide_width - mm(24), prs.slide_height - mm(12), mm(18), mm(7))
    tf = tb.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = text
    apply_run_font(r, size_pt=9, color=ACCENT_SUB)


def layout_title_slide(prs: Presentation, idx: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_base_bg(slide, prs, gradient=False)

    t = add_textbox(slide, mm(24), mm(58), mm(249), mm(28))
    tf = t.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = f"{{{{title_{idx}}}}}"
    apply_run_font(r, size_pt=40, color=TEXT_WHITE, bold=True)

    s = add_textbox(slide, mm(24), mm(92), mm(249), mm(16))
    sf = s.text_frame
    sf.clear()
    sp = sf.paragraphs[0]
    sp.alignment = PP_ALIGN.CENTER
    sr = sp.add_run()
    sr.text = f"{{{{subtitle_{idx}}}}}"
    apply_run_font(sr, size_pt=18, color=TEXT_LIGHT)

    d = add_textbox(slide, mm(24), mm(118), mm(249), mm(10))
    df = d.text_frame
    df.clear()
    dp = df.paragraphs[0]
    dp.alignment = PP_ALIGN.CENTER
    dr = dp.add_run()
    dr.text = f"{{{{date_{idx}}}}}"
    apply_run_font(dr, size_pt=12, color=ACCENT_MAIN)

    add_slide_number(slide, prs, f"{{{{slide_no_{idx}}}}}")


def layout_agenda_slide(prs: Presentation, idx: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_base_bg(slide, prs, gradient=False)

    h = add_textbox(slide, mm(18), mm(14), mm(220), mm(12))
    hf = h.text_frame
    hf.clear()
    hp = hf.paragraphs[0]
    hr = hp.add_run()
    hr.text = f"{{{{title_{idx}}}}}"
    apply_run_font(hr, size_pt=24, color=TEXT_WHITE, bold=True)

    card_x, card_y, card_w, card_h = mm(18), mm(50), mm(261), mm(85)
    add_rect(slide, card_x, card_y, card_w, card_h, fill=CARD_BG, line=ACCENT_MAIN, rounded=True)

    base_y = i(card_y + mm(22))
    step_w = mm(58)
    gap = mm(5)
    x0 = i(card_x + mm(8))

    for n in range(4):
        sx = i(x0 + n * (step_w + gap))

        c = add_rect(slide, sx, base_y, mm(9), mm(9), fill=ACCENT_SUB, rounded=True)
        ctf = c.text_frame
        ctf.clear()
        cp = ctf.paragraphs[0]
        cp.alignment = PP_ALIGN.CENTER
        cr = cp.add_run()
        cr.text = str(n + 1)
        apply_run_font(cr, size_pt=10, color=BG_DARK, bold=True)

        tb = add_textbox(slide, sx + mm(11), base_y - mm(1), step_w - mm(11), mm(16))
        tf = tb.text_frame
        tf.clear()
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = f"{{{{step{n+1}_{idx}}}}}"
        apply_run_font(r, size_pt=12, color=TEXT_LIGHT)
        p.line_spacing = 1.2

        if n < 3:
            arr = add_textbox(slide, sx + step_w + mm(1), base_y + mm(1), mm(4), mm(8))
            atf = arr.text_frame
            atf.clear()
            ap = atf.paragraphs[0]
            ap.alignment = PP_ALIGN.CENTER
            ar = ap.add_run()
            ar.text = "→"
            apply_run_font(ar, size_pt=15, color=ACCENT_MAIN, bold=True)

    add_slide_number(slide, prs, f"{{{{slide_no_{idx}}}}}")


def layout_section_slide(prs: Presentation, idx: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_base_bg(slide, prs, gradient=True)

    t = add_textbox(slide, mm(22), mm(80), mm(253), mm(40))
    tf = t.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = f"{{{{title_{idx}}}}}"
    apply_run_font(r, size_pt=36, color=TEXT_WHITE, bold=True)

    add_slide_number(slide, prs, f"{{{{slide_no_{idx}}}}}")


def layout_content_slide(prs: Presentation, idx: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_base_bg(slide, prs, gradient=False)

    h = add_textbox(slide, mm(18), mm(14), mm(250), mm(12))
    hf = h.text_frame
    hf.clear()
    hp = hf.paragraphs[0]
    hr = hp.add_run()
    hr.text = f"{{{{title_{idx}}}}}"
    apply_run_font(hr, size_pt=24, color=TEXT_WHITE, bold=True)

    card = add_rect(slide, mm(18), mm(34), mm(261), mm(145), fill=CARD_BG, line=ACCENT_MAIN, rounded=True)
    card.line.width = Pt(1.2)

    b = add_textbox(slide, mm(26), mm(44), mm(245), mm(128))
    bf = b.text_frame
    bf.clear()
    bf.word_wrap = True
    p = bf.paragraphs[0]
    r = p.add_run()
    r.text = f"{{{{body_{idx}}}}}"
    apply_run_font(r, size_pt=16, color=TEXT_LIGHT)
    p.line_spacing = 1.35

    add_slide_number(slide, prs, f"{{{{slide_no_{idx}}}}}")


def layout_split_slide(prs: Presentation, idx: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_base_bg(slide, prs, gradient=False)

    h = add_textbox(slide, mm(18), mm(14), mm(250), mm(12))
    hf = h.text_frame
    hf.clear()
    hp = hf.paragraphs[0]
    hr = hp.add_run()
    hr.text = f"{{{{title_{idx}}}}}"
    apply_run_font(hr, size_pt=24, color=TEXT_WHITE, bold=True)

    add_rect(slide, mm(18), mm(34), mm(126), mm(145), fill=CARD_BG, line=ACCENT_MAIN, rounded=True)
    lt = add_textbox(slide, mm(24), mm(42), mm(114), mm(12))
    ltf = lt.text_frame
    ltf.clear()
    lp = ltf.paragraphs[0]
    lr = lp.add_run()
    lr.text = f"{{{{left_title_{idx}}}}}"
    apply_run_font(lr, size_pt=18, color=ACCENT_MAIN, bold=True)

    lb = add_textbox(slide, mm(24), mm(56), mm(114), mm(116))
    lbf = lb.text_frame
    lbf.clear()
    lbf.word_wrap = True
    lpb = lbf.paragraphs[0]
    lrb = lpb.add_run()
    lrb.text = f"{{{{left_body_{idx}}}}}"
    apply_run_font(lrb, size_pt=14, color=TEXT_LIGHT)
    lpb.line_spacing = 1.3

    add_rect(slide, mm(153), mm(34), mm(126), mm(145), fill=CARD_BG, line=ACCENT_MAIN, rounded=True)
    rt = add_textbox(slide, mm(159), mm(42), mm(114), mm(12))
    rtf = rt.text_frame
    rtf.clear()
    rp = rtf.paragraphs[0]
    rr = rp.add_run()
    rr.text = f"{{{{right_title_{idx}}}}}"
    apply_run_font(rr, size_pt=18, color=ACCENT_MAIN, bold=True)

    rb = add_textbox(slide, mm(159), mm(56), mm(114), mm(116))
    rbf = rb.text_frame
    rbf.clear()
    rbf.word_wrap = True
    rpb = rbf.paragraphs[0]
    rrb = rpb.add_run()
    rrb.text = f"{{{{right_body_{idx}}}}}"
    apply_run_font(rrb, size_pt=14, color=TEXT_LIGHT)
    rpb.line_spacing = 1.3

    add_slide_number(slide, prs, f"{{{{slide_no_{idx}}}}}")


def layout_summary_slide(prs: Presentation, idx: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_base_bg(slide, prs, gradient=False)

    h = add_textbox(slide, mm(18), mm(14), mm(250), mm(12))
    hf = h.text_frame
    hf.clear()
    hp = hf.paragraphs[0]
    hr = hp.add_run()
    hr.text = f"{{{{title_{idx}}}}}"
    apply_run_font(hr, size_pt=24, color=TEXT_WHITE, bold=True)

    card = add_rect(slide, mm(18), mm(34), mm(261), mm(120), fill=CARD_BG, line=ACCENT_MAIN, rounded=True)
    card.line.width = Pt(1.2)

    b = add_textbox(slide, mm(26), mm(44), mm(245), mm(96))
    bf = b.text_frame
    bf.clear()
    bf.word_wrap = True
    p = bf.paragraphs[0]
    r = p.add_run()
    r.text = f"{{{{body_{idx}}}}}"
    apply_run_font(r, size_pt=16, color=TEXT_LIGHT)
    p.line_spacing = 1.35

    c = add_textbox(slide, mm(18), mm(160), mm(261), mm(20))
    cf = c.text_frame
    cf.clear()
    cp = cf.paragraphs[0]
    cp.alignment = PP_ALIGN.CENTER
    cr = cp.add_run()
    cr.text = f"{{{{contact_{idx}}}}} / {{{{organizer_{idx}}}}}"
    apply_run_font(cr, size_pt=12, color=ACCENT_SUB, bold=True)

    add_slide_number(slide, prs, f"{{{{slide_no_{idx}}}}}")


def create_slides_template(output_path: str | Path, slides: List[Dict[str, Any]]) -> Path:
    prs = Presentation()
    prs.slide_width = mm(297)   # A4 landscape
    prs.slide_height = mm(210)

    for idx, s in enumerate(slides, start=1):
        t = s.get("type")
        if t == "title":
            layout_title_slide(prs, idx)
        elif t == "agenda":
            layout_agenda_slide(prs, idx)
        elif t == "section":
            layout_section_slide(prs, idx)
        elif t == "content":
            layout_content_slide(prs, idx)
        elif t == "split":
            layout_split_slide(prs, idx)
        elif t == "summary":
            layout_summary_slide(prs, idx)
        else:
            raise ValueError(f"Unknown slide type: {t}")

    out = Path(output_path).expanduser()
    out.parent.mkdir(parents=True, exist_ok=True)

    try:
        prs.save(str(out))
    except Exception:
        traceback.print_exc()
        raise

    return out


if __name__ == "__main__":
    slide_defs = [
        {"type": "title"},
        {"type": "agenda"},
        {"type": "section"},
        {"type": "content"},
        {"type": "content"},
        {"type": "section"},
        {"type": "split"},
        {"type": "content"},
        {"type": "section"},
        {"type": "content"},
        {"type": "content"},
        {"type": "content"},
        {"type": "summary"},
    ]

    template_out = "/home/tomoyuki/.openclaw/workspace/output/komagane_ai_share_slides_template.pptx"
    print(create_slides_template(template_out, slide_defs))
