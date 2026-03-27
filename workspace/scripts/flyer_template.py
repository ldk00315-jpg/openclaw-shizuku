from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path
from typing import Any, Dict, List

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


def mm(v: float) -> float:
    return Inches(v / 25.4)


@dataclass(frozen=True)
class Palette:
    primary: RGBColor = RGBColor(0xE8, 0x82, 0x3A)      # #E8823A
    secondary: RGBColor = RGBColor(0x6B, 0xB7, 0x7B)    # #6BB77B
    body: RGBColor = RGBColor(0x33, 0x33, 0x33)         # #333333
    white: RGBColor = RGBColor(0xFF, 0xFF, 0xFF)        # #FFFFFF
    bg: RGBColor = RGBColor(0xFF, 0xF8, 0xF0)           # #FFF8F0
    divider: RGBColor = RGBColor(0xF5, 0xC5, 0x9F)      # #F5C59F
    schedule_bg: RGBColor = RGBColor(0xEA, 0xF5, 0xEC)  # #EAF5EC
    footer_bg: RGBColor = RGBColor(0x44, 0x44, 0x44)    # #444444


DEFAULT_PLACEHOLDERS: Dict[str, Any] = {
    "title": "{{title}}",
    "subtitle": "{{subtitle}}",
    "about_text": "{{about_text}}",
    "schedule_steps": ["{{step1}}", "{{step2}}", "{{step3}}", "{{step4}}"],
    "goals": ["{{goal1}}", "{{goal2}}", "{{goal3}}"],
    "rules": ["{{rule1}}", "{{rule2}}", "{{rule3}}"],
    "organizer": "{{organizer}}",
    "cooperation": "{{cooperation}}",
    "contact": "{{contact}}",
    "date_location": "{{date_location}}",
}


def _set_font(run, name: str, size_pt: float, color: RGBColor, bold: bool = False):
    run.font.name = name
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = color


def _add_rect(slide, x, y, w, h, fill: RGBColor | None = None, line: RGBColor | None = None, radius: bool = False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, x, y, w, h)
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    return shp


def _add_section_title(slide, text: str, x, y, w, color: RGBColor, font_name: str = "Noto Sans JP"):
    tb = slide.shapes.add_textbox(x, y, w, mm(8))
    tf = tb.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    _set_font(run, font_name, 14, color, bold=True)
    return tb


def _add_bullets(slide, items: List[str], x, y, w, h, dot_color: RGBColor, font_size: float = 11, line_spacing: float = 1.3):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.clear()
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"● {item}"
        p.line_spacing = line_spacing
        p.space_after = Pt(2)
        if p.runs:
            _set_font(p.runs[0], "Noto Sans JP", font_size, dot_color, bold=False)
            # color bullet+text uniformly for simplicity, then recolor text body if needed
            p.runs[0].font.color.rgb = dot_color
    return tb


def create_flyer(content: Dict[str, Any], out_path: str | Path):
    c = {**DEFAULT_PLACEHOLDERS, **(content or {})}
    palette = Palette()

    prs = Presentation()
    prs.slide_width = mm(210)
    prs.slide_height = mm(297)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # page + margins
    margin = mm(12)
    content_x = margin
    content_w = mm(210 - 24)
    gap = mm(1.6)

    # full background
    _add_rect(slide, 0, 0, prs.slide_width, prs.slide_height, fill=palette.bg)

    y = margin

    # 1) Header (50mm)
    header_h = mm(50)
    header = _add_rect(slide, content_x, y, content_w, header_h, fill=palette.primary)
    htf = header.text_frame
    htf.clear()
    p1 = htf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run()
    r1.text = c["title"]
    _set_font(r1, "Noto Sans JP", 28, palette.white, bold=True)

    p2 = htf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = c["subtitle"]
    _set_font(r2, "Noto Sans JP", 16, palette.white, bold=False)

    y += header_h + gap

    # divider
    slide.shapes.add_connector(1, content_x, y - gap / 2, content_x + content_w, y - gap / 2).line.color.rgb = palette.divider

    # 2) About (55mm)
    about_h = mm(55)
    _add_rect(slide, content_x, y, content_w, about_h, fill=palette.bg)
    _add_section_title(slide, "こんな会です", content_x + mm(8), y + mm(4), content_w - mm(8), palette.primary)
    # green round icon
    _add_rect(slide, content_x + mm(2), y + mm(4), mm(4), mm(4), fill=palette.secondary, radius=True)

    about_tb = slide.shapes.add_textbox(content_x + mm(8), y + mm(14), content_w - mm(10), about_h - mm(16))
    about_tf = about_tb.text_frame
    about_tf.clear()
    for idx, line in enumerate(str(c["about_text"]).split("\n")):
        p = about_tf.paragraphs[0] if idx == 0 else about_tf.add_paragraph()
        run = p.add_run()
        run.text = line
        _set_font(run, "Noto Sans JP", 11, palette.body, bold=False)
        p.line_spacing = 1.3

    y += about_h + gap
    slide.shapes.add_connector(1, content_x, y - gap / 2, content_x + content_w, y - gap / 2).line.color.rgb = palette.divider

    # 3) Schedule (50mm)
    sch_h = mm(50)
    _add_rect(slide, content_x, y, content_w, sch_h, fill=palette.schedule_bg)
    _add_section_title(slide, "初回のイメージ", content_x + mm(2), y + mm(3), content_w - mm(4), palette.primary)

    steps = list(c.get("schedule_steps", []))[:4]
    while len(steps) < 4:
        steps.append("")

    base_y = y + mm(15)
    col_w = mm((210 - 24) / 4 - 2)
    start_x = content_x + mm(1.5)
    for i, s in enumerate(steps):
        sx = start_x + i * (col_w + mm(1.5))
        # green circle + number
        circle = _add_rect(slide, sx, base_y, mm(8), mm(8), fill=palette.secondary, radius=True)
        ctf = circle.text_frame
        ctf.clear()
        cp = ctf.paragraphs[0]
        cp.alignment = PP_ALIGN.CENTER
        cr = cp.add_run()
        cr.text = str(i + 1)
        _set_font(cr, "Noto Sans JP", 11, palette.white, bold=True)

        st = slide.shapes.add_textbox(sx + mm(9.5), base_y - mm(0.5), col_w - mm(10), mm(12))
        stf = st.text_frame
        stf.clear()
        sp = stf.paragraphs[0]
        sr = sp.add_run()
        sr.text = s
        _set_font(sr, "Noto Sans JP", 10.5, palette.body)
        sp.line_spacing = 1.3

        if i < 3:
            arr = slide.shapes.add_textbox(sx + col_w + mm(0.4), base_y + mm(1.2), mm(1.2), mm(6))
            atf = arr.text_frame
            atf.clear()
            ap = atf.paragraphs[0]
            ap.alignment = PP_ALIGN.CENTER
            ar = ap.add_run()
            ar.text = "→"
            _set_font(ar, "Noto Sans JP", 14, palette.secondary, bold=True)

    y += sch_h + gap
    slide.shapes.add_connector(1, content_x, y - gap / 2, content_x + content_w, y - gap / 2).line.color.rgb = palette.divider

    # 4) Goals (45mm)
    goals_h = mm(45)
    _add_rect(slide, content_x, y, content_w, goals_h, fill=palette.bg)
    _add_section_title(slide, "この会で目指すこと", content_x + mm(2), y + mm(4), content_w - mm(4), palette.primary)
    _add_bullets(
        slide,
        list(c.get("goals", []))[:4],
        content_x + mm(3),
        y + mm(14),
        content_w - mm(6),
        goals_h - mm(16),
        dot_color=palette.secondary,
        font_size=11,
        line_spacing=1.3,
    )

    y += goals_h + gap
    slide.shapes.add_connector(1, content_x, y - gap / 2, content_x + content_w, y - gap / 2).line.color.rgb = palette.divider

    # 5) Rules (30mm)
    rules_h = mm(30)
    _add_rect(slide, content_x, y, content_w, rules_h, fill=palette.bg)
    _add_section_title(slide, "場のルール", content_x + mm(2), y + mm(3), content_w - mm(4), palette.primary)
    _add_bullets(
        slide,
        list(c.get("rules", []))[:4],
        content_x + mm(3),
        y + mm(11),
        content_w - mm(6),
        rules_h - mm(12),
        dot_color=palette.primary,
        font_size=10,
        line_spacing=1.3,
    )

    y += rules_h + gap
    slide.shapes.add_connector(1, content_x, y - gap / 2, content_x + content_w, y - gap / 2).line.color.rgb = palette.divider

    # 6) Footer (35mm)
    footer_h = mm(35)
    _add_rect(slide, content_x, y, content_w, footer_h, fill=palette.footer_bg)

    left = slide.shapes.add_textbox(content_x + mm(4), y + mm(5), content_w * 0.48, footer_h - mm(8))
    ltf = left.text_frame
    ltf.clear()
    lp1 = ltf.paragraphs[0]
    lr1 = lp1.add_run()
    lr1.text = f"主催：{c['organizer']}"
    _set_font(lr1, "Noto Sans JP", 9, palette.white)
    lp2 = ltf.add_paragraph()
    lr2 = lp2.add_run()
    lr2.text = f"協力：{c['cooperation']}"
    _set_font(lr2, "Noto Sans JP", 9, palette.white)

    right = slide.shapes.add_textbox(content_x + content_w * 0.52, y + mm(5), content_w * 0.44, footer_h - mm(8))
    rtf = right.text_frame
    rtf.clear()
    rp1 = rtf.paragraphs[0]
    rp1.alignment = PP_ALIGN.RIGHT
    rr1 = rp1.add_run()
    rr1.text = c["date_location"]
    _set_font(rr1, "Noto Sans JP", 9, palette.white)
    rp2 = rtf.add_paragraph()
    rp2.alignment = PP_ALIGN.RIGHT
    rr2 = rp2.add_run()
    rr2.text = c["contact"]
    _set_font(rr2, "Noto Sans JP", 9, palette.white)

    out_path = Path(out_path).expanduser()
    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    return out_path


def create_template(out_path: str | Path):
    return create_flyer(DEFAULT_PLACEHOLDERS, out_path)
