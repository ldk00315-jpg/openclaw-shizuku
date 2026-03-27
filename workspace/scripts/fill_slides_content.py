#!/home/tomoyuki/.openclaw/workspace/.venv-pptx/bin/python
from __future__ import annotations

from pathlib import Path
from typing import Dict, List, Any

from pptx import Presentation


def build_replacements(slides: List[Dict[str, Any]]) -> Dict[str, str]:
    repl: Dict[str, str] = {}

    for i, s in enumerate(slides, start=1):
        repl[f"{{{{slide_no_{i}}}}}"] = str(i)

        if "title" in s:
            repl[f"{{{{title_{i}}}}}"] = s.get("title", "")

        st = s.get("type")
        if st == "title":
            repl[f"{{{{subtitle_{i}}}}}"] = s.get("subtitle", "")
            repl[f"{{{{date_{i}}}}}"] = s.get("date", "")

        elif st == "agenda":
            steps = s.get("steps", ["", "", "", ""])
            for n in range(4):
                repl[f"{{{{step{n+1}_{i}}}}}"] = steps[n] if n < len(steps) else ""

        elif st == "content":
            repl[f"{{{{body_{i}}}}}"] = s.get("body", "")

        elif st == "split":
            repl[f"{{{{left_title_{i}}}}}"] = s.get("left_title", "")
            repl[f"{{{{left_body_{i}}}}}"] = s.get("left_body", "")
            repl[f"{{{{right_title_{i}}}}}"] = s.get("right_title", "")
            repl[f"{{{{right_body_{i}}}}}"] = s.get("right_body", "")

        elif st == "summary":
            repl[f"{{{{body_{i}}}}}"] = s.get("body", "")
            repl[f"{{{{contact_{i}}}}}"] = s.get("contact", "")
            repl[f"{{{{organizer_{i}}}}}"] = s.get("organizer", "")

    return repl


def fill_slides(template_path: str | Path, output_path: str | Path, slides: List[Dict[str, Any]]) -> Path:
    prs = Presentation(str(Path(template_path).expanduser()))
    repl = build_replacements(slides)

    for slide in prs.slides:
        for shape in slide.shapes:
            if not hasattr(shape, "text_frame") or not shape.text_frame:
                continue
            tf = shape.text_frame
            for p in tf.paragraphs:
                for run in p.runs:
                    txt = run.text
                    for key, val in repl.items():
                        if key in txt:
                            txt = txt.replace(key, val)
                    run.text = txt

    out = Path(output_path).expanduser()
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    return out


if __name__ == "__main__":
    slides = [
        {"type": "title", "title": "駒ヶ根AI実践シェア会", "subtitle": "成果より、まず「やってみた」を持ち寄る会", "date": "第1回：2026年4月（日程調整中）"},
        {"type": "agenda", "title": "今日の流れ", "steps": ["チェックイン\n自己紹介", "トーク\n体験共有", "グループ\nワーク", "まとめ\n次回予告"]},
        {"type": "section", "title": "AIチャットとは"},
        {"type": "content", "title": "AIチャットの基本と使い方", "body": "・質問を入力すると、AIが回答を生成\n・ChatGPT、Gemini、Claude などが代表的\n・スマホでもPCでも使える\n・「完璧な質問」じゃなくてOK、気軽に聞いてみよう"},
        {"type": "content", "title": "身近なAI活用事例", "body": "・メール文面の作成・添削\n・議事録の要約\n・チラシやSNS投稿の文案\n・翻訳・言い換え\n・アイデア出しの壁打ち相手\n※このチラシもAIエージェント（しずく）が作成しました"},
        {"type": "section", "title": "AIエージェントとは"},
        {"type": "split", "title": "Web型 vs ローカル型", "left_title": "Web型", "left_body": "・ブラウザで動作\n・すぐ使える\n・例：ChatGPT、Gemini\n・インターネット接続が必要", "right_title": "ローカル型", "right_body": "・自分のPCで動作\n・データが外に出ない\n・例：OpenClaw、Ollama\n・初期設定が必要"},
        {"type": "content", "title": "AIエージェントで出来ること", "body": "・ファイル作成（文書、スライド、表）\n・情報の検索・整理\n・定型作業の自動化\n・スケジュール管理の補助\n・コードの生成・修正"},
        {"type": "section", "title": "OpenClaw実践\n（初級編）"},
        {"type": "content", "title": "OpenClawデモ — このチラシができるまで", "body": "1. Telegramでしずく（AIエージェント）に指示\n2. python-pptxでPPTX生成スクリプトを作成\n3. テンプレート生成 → コンテンツ注入の2段階\n4. 検証 → 修正 → 再生成のサイクル\n5. 完成品をTelegram＋メールで自動送信"},
        {"type": "content", "title": "「AIでこれやってみよう」ワーク", "body": "スマホでAIチャットを開いてみましょう！\n\nお題：\n・自分の仕事の悩みをAIに相談してみよう\n・来週のメール文面を下書きしてもらおう\n・好きな料理のアレンジレシピを聞いてみよう\n\n5分間、自由に試してみてください"},
        {"type": "content", "title": "AIとの付き合い方", "body": "・AIの回答は必ずしも正確ではない（ハルシネーション）\n・個人情報・機密情報は入力しない\n・著作権のある文章をそのままコピペしない\n・「便利なアシスタント」として使う、丸投げしない\n・知っておけば怖くない、まず使ってみよう"},
        {"type": "summary", "title": "まとめ・次回予告", "body": "・AIは「まず触ってみる」が一番の近道\n・失敗OK、試行錯誤を楽しもう\n・次回も体験シェアをお待ちしています", "contact": "お問い合わせ：valuegarage@gmail.com", "organizer": "主催：南信AI実践会 協力：駒ヶ根商工会議所（予定）"},
    ]

    template = "/home/tomoyuki/.openclaw/workspace/output/komagane_ai_share_slides_template.pptx"
    output = "/home/tomoyuki/.openclaw/workspace/output/komagane_ai_share_slides.pptx"
    print(fill_slides(template, output, slides))
