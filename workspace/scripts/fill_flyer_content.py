from __future__ import annotations

from pathlib import Path

from pptx import Presentation


def fill_template(template_path: Path, output_path: Path, content: dict) -> Path:
    prs = Presentation(str(template_path))

    replacements = {
        "{{title}}": content["title"],
        "{{subtitle}}": content["subtitle"],
        "{{about_text}}": content["about_text"],
        "{{step1}}": content["schedule_steps"][0],
        "{{step2}}": content["schedule_steps"][1],
        "{{step3}}": content["schedule_steps"][2],
        "{{step4}}": content["schedule_steps"][3],
        "{{goal1}}": content["goals"][0],
        "{{goal2}}": content["goals"][1],
        "{{goal3}}": content["goals"][2],
        "{{goal4}}": content["goals"][3],
        "{{rule1}}": content["rules"][0],
        "{{rule2}}": content["rules"][1],
        "{{rule3}}": content["rules"][2],
        "{{rule4}}": content["rules"][3],
        "{{organizer}}": content["organizer"],
        "{{cooperation}}": content["cooperation"],
        "{{contact}}": content["contact"],
        "{{date_location}}": content["date_location"],
    }

    for slide in prs.slides:
        for shape in slide.shapes:
            if not hasattr(shape, "text_frame") or not shape.text_frame:
                continue
            tf = shape.text_frame
            for paragraph in tf.paragraphs:
                for run in paragraph.runs:
                    text = run.text
                    for old, new in replacements.items():
                        if old in text:
                            text = text.replace(old, new)
                    run.text = text

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    return output_path


if __name__ == "__main__":
    content = {
        "title": "駒ヶ根AI実践シェア会",
        "subtitle": "成果より、まず「やってみた」を持ち寄る会",
        "about_text": "AIを使ってみた体験を気軽にシェアする場です。\n専門知識は不要。「こんなの試してみた」「うまくいかなかった」\nそんな話が一番盛り上がります。個人・小規模事業者の方、大歓迎。",
        "schedule_steps": ["チェックイン\n自己紹介", "トーク\n体験共有", "グループ\nワーク", "まとめ\n次回予告"],
        "goals": ["AIを身近に感じる", "失敗も含めて共有できる安心の場", "地域のAI仲間をつくる", "小さな一歩を応援し合う"],
        "rules": ["否定しない、まず聴く", "専門用語より体験談", "SNS投稿は本人の許可を得て", "楽しむことが最優先"],
        "organizer": "南信AI実践会",
        "cooperation": "駒ヶ根商工会議所（予定）",
        "contact": "お問い合わせ：valuegarage@gmail.com",
        "date_location": "第1回：2026年4月（日程調整中）\n会場：駒ヶ根市内（調整中）",
    }

    template = Path("/home/tomoyuki/.openclaw/workspace/templates/flyer_a4_portrait.pptx")
    output = Path("/home/tomoyuki/.openclaw/workspace/output/komagane_ai_share_flyer.pptx")
    result = fill_template(template, output, content)
    print(result)
