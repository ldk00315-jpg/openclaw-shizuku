#!/home/tomoyuki/.openclaw/workspace/.venv-pptx/bin/python
from __future__ import annotations

import traceback
from pathlib import Path
from typing import Any, Dict

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt

FONT_FAMILY = "Noto Sans CJK JP"


def mm(v: float) -> int:
    return int(Inches(v / 25.4))


def as_int(v) -> int:
    return int(v)


def apply_run_font(run, *, size_pt: float, color: RGBColor, bold: bool = False, font_name: str = FONT_FAMILY) -> None:
    run.font.name = font_name
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = color

    r = run._r
    rPr = r.get_or_add_rPr()
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = OxmlElement("a:ea")
        rPr.append(ea)
    ea.set("typeface", font_name)


def create_flyer_template(template_path: str | Path) -> Path:
    prs = Presentation()
    prs.slide_width = mm(210)
    prs.slide_height = mm(297)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    PRIMARY = RGBColor(0xE8, 0x82, 0x3A)
    SECONDARY = RGBColor(0x6B, 0xB7, 0x7B)
    BODY = RGBColor(0x33, 0x33, 0x33)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    BG = RGBColor(0xFF, 0xF8, 0xF0)
    DIVIDER = RGBColor(0xF5, 0xC5, 0x9F)
    SCHEDULE_BG = RGBColor(0xEA, 0xF5, 0xEC)
    FOOTER_BG = RGBColor(0x44, 0x44, 0x44)

    data: Dict[str, Any] = {
        "title": "{{title}}",
        "subtitle": "{{subtitle}}",
        "about_text": "{{about_text}}",
        "schedule_steps": ["{{step1}}", "{{step2}}", "{{step3}}", "{{step4}}"],
        "goals": ["{{goal1}}", "{{goal2}}", "{{goal3}}", "{{goal4}}"],
        "rules": ["{{rule1}}", "{{rule2}}", "{{rule3}}", "{{rule4}}"],
        "organizer": "{{organizer}}",
        "cooperation": "{{cooperation}}",
        "contact": "{{contact}}",
        "date_location": "{{date_location}}",
    }

    margin = mm(12)
    content_x = margin
    content_w = mm(210 - 24)
    gap = mm(2)

    def add_textbox(x, y, w, h):
        return slide.shapes.add_textbox(as_int(x), as_int(y), as_int(w), as_int(h))

    def add_rect(x, y, w, h, fill=None, line=None, rounded=False):
        shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
        shp = slide.shapes.add_shape(shape_type, as_int(x), as_int(y), as_int(w), as_int(h))
        if fill is None:
            shp.fill.background()
        else:
            shp.fill.solid()
            shp.fill.fore_color.rgb = fill
        if line is None:
            shp.line.fill.background()
        else:
            shp.line.color.rgb = line
        return shp

    def add_divider(y_pos):
        line = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            as_int(content_x),
            as_int(y_pos),
            as_int(content_x + content_w),
            as_int(y_pos),
        )
        line.line.color.rgb = DIVIDER

    add_rect(0, 0, prs.slide_width, prs.slide_height, fill=BG)

    y = margin

    header_h = mm(50)
    header = add_rect(content_x, y, content_w, header_h, fill=PRIMARY)
    tf = header.text_frame
    tf.clear()

    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run()
    r1.text = data["title"]
    apply_run_font(r1, size_pt=28, color=WHITE, bold=True)

    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = data["subtitle"]
    apply_run_font(r2, size_pt=16, color=WHITE)

    y = as_int(y + header_h + gap)
    add_divider(as_int(y - gap / 2))

    about_h = mm(55)
    add_rect(content_x, y, content_w, about_h, fill=BG)
    icon = add_rect(content_x + mm(2), y + mm(5), mm(4.5), mm(4.5), fill=SECONDARY, rounded=True)
    icon.line.fill.background()

    title_tb = add_textbox(content_x + mm(8), y + mm(4), content_w - mm(8), mm(8))
    ttf = title_tb.text_frame
    ttf.clear()
    tp = ttf.paragraphs[0]
    tr = tp.add_run()
    tr.text = "こんな会です"
    apply_run_font(tr, size_pt=14, color=PRIMARY, bold=True)

    body_tb = add_textbox(content_x + mm(8), y + mm(14), content_w - mm(10), about_h - mm(16))
    btf = body_tb.text_frame
    btf.clear()
    bp = btf.paragraphs[0]
    br = bp.add_run()
    br.text = data["about_text"]
    apply_run_font(br, size_pt=11, color=BODY)
    bp.line_spacing = 1.3

    y = as_int(y + about_h + gap)
    add_divider(as_int(y - gap / 2))

    sch_h = mm(50)
    add_rect(content_x, y, content_w, sch_h, fill=SCHEDULE_BG)

    sch_title_tb = add_textbox(content_x + mm(2), y + mm(3), content_w - mm(4), mm(8))
    stf = sch_title_tb.text_frame
    stf.clear()
    sp = stf.paragraphs[0]
    sr = sp.add_run()
    sr.text = "初回のイメージ"
    apply_run_font(sr, size_pt=14, color=PRIMARY, bold=True)

    steps = list(data["schedule_steps"])[:4]
    while len(steps) < 4:
        steps.append("")

    base_y = as_int(y + mm(15))
    step_w = mm(40)
    step_gap = mm(3)
    x0 = as_int(content_x + mm(2))

    for i, txt in enumerate(steps):
        sx = as_int(x0 + i * (step_w + step_gap))

        circ = add_rect(sx, base_y, mm(8), mm(8), fill=SECONDARY, rounded=True)
        ctf = circ.text_frame
        ctf.clear()
        cp = ctf.paragraphs[0]
        cp.alignment = PP_ALIGN.CENTER
        cr = cp.add_run()
        cr.text = str(i + 1)
        apply_run_font(cr, size_pt=11, color=WHITE, bold=True)

        tb = add_textbox(sx + mm(9.5), base_y - mm(0.5), step_w - mm(9.5), mm(13))
        tbf = tb.text_frame
        tbf.clear()
        tbf.word_wrap = True
        tp = tbf.paragraphs[0]
        tr = tp.add_run()
        tr.text = txt
        apply_run_font(tr, size_pt=10, color=BODY)
        tp.line_spacing = 1.3

        if i < 3:
            arr = add_textbox(sx + step_w + mm(0.3), base_y + mm(1.2), mm(2), mm(6))
            atf = arr.text_frame
            atf.clear()
            ap = atf.paragraphs[0]
            ap.alignment = PP_ALIGN.CENTER
            ar = ap.add_run()
            ar.text = "→"
            apply_run_font(ar, size_pt=14, color=SECONDARY, bold=True)

    y = as_int(y + sch_h + gap)
    add_divider(as_int(y - gap / 2))

    goals_h = mm(45)
    add_rect(content_x, y, content_w, goals_h, fill=BG)

    gtitle_tb = add_textbox(content_x + mm(2), y + mm(4), content_w - mm(4), mm(8))
    gtf = gtitle_tb.text_frame
    gtf.clear()
    gp = gtf.paragraphs[0]
    gr = gp.add_run()
    gr.text = "この会で目指すこと"
    apply_run_font(gr, size_pt=14, color=PRIMARY, bold=True)

    goals_tb = add_textbox(content_x + mm(3), y + mm(13), content_w - mm(6), goals_h - mm(15))
    gbtf = goals_tb.text_frame
    gbtf.clear()
    for idx, item in enumerate(data["goals"]):
        p = gbtf.paragraphs[0] if idx == 0 else gbtf.add_paragraph()
        r = p.add_run()
        r.text = f"● {item}"
        apply_run_font(r, size_pt=11, color=SECONDARY)
        p.line_spacing = 1.3

    y = as_int(y + goals_h + gap)
    add_divider(as_int(y - gap / 2))

    rules_h = mm(35)
    add_rect(content_x, y, content_w, rules_h, fill=BG)

    rtitle_tb = add_textbox(content_x + mm(2), y + mm(3), content_w - mm(4), mm(7))
    rtf = rtitle_tb.text_frame
    rtf.clear()
    rp = rtf.paragraphs[0]
    rr = rp.add_run()
    rr.text = "場のルール"
    apply_run_font(rr, size_pt=14, color=PRIMARY, bold=True)

    rules_tb = add_textbox(content_x + mm(3), y + mm(10), content_w - mm(6), rules_h - mm(12))
    rbtf = rules_tb.text_frame
    rbtf.clear()
    for idx, item in enumerate(data["rules"]):
        p = rbtf.paragraphs[0] if idx == 0 else rbtf.add_paragraph()
        r = p.add_run()
        r.text = f"● {item}"
        apply_run_font(r, size_pt=10, color=PRIMARY)
        p.line_spacing = 1.3

    y = as_int(y + rules_h + gap)
    add_divider(as_int(y - gap / 2))

    footer_h = mm(35)
    add_rect(content_x, y, content_w, footer_h, fill=FOOTER_BG)

    left_tb = add_textbox(content_x + mm(4), y + mm(5), int(content_w * 0.48), footer_h - mm(8))
    ltf = left_tb.text_frame
    ltf.clear()

    lp1 = ltf.paragraphs[0]
    lr1 = lp1.add_run()
    lr1.text = f"主催：{data['organizer']}"
    apply_run_font(lr1, size_pt=9, color=WHITE)

    lp2 = ltf.add_paragraph()
    lr2 = lp2.add_run()
    lr2.text = f"協力：{data['cooperation']}"
    apply_run_font(lr2, size_pt=9, color=WHITE)

    right_tb = add_textbox(content_x + int(content_w * 0.52), y + mm(5), int(content_w * 0.44), footer_h - mm(8))
    rtf = right_tb.text_frame
    rtf.clear()

    rp1 = rtf.paragraphs[0]
    rp1.alignment = PP_ALIGN.RIGHT
    rr1 = rp1.add_run()
    rr1.text = data["date_location"]
    apply_run_font(rr1, size_pt=9, color=WHITE)

    rp2 = rtf.add_paragraph()
    rp2.alignment = PP_ALIGN.RIGHT
    rr2 = rp2.add_run()
    rr2.text = data["contact"]
    apply_run_font(rr2, size_pt=9, color=WHITE)

    template_path = Path(template_path).expanduser()
    template_path.parent.mkdir(parents=True, exist_ok=True)

    try:
        prs.save(str(template_path))
    except Exception:
        traceback.print_exc()
        raise

    return template_path


if __name__ == "__main__":
    out = Path("~/.openclaw/workspace/templates/flyer_a4_portrait.pptx").expanduser()
    path = create_flyer_template(out)
    print(path)
