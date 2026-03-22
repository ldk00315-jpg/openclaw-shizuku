#!/usr/bin/env python3
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

OUT = "/home/tomoyuki/.openclaw/workspace/notes/ai-agent-study-2026-03-22.pptx"

prs = Presentation()

# theme colors
NAVY = RGBColor(30, 58, 138)
GRAY = RGBColor(75, 85, 99)
ACCENT = RGBColor(14, 116, 144)
BG = RGBColor(248, 250, 252)


def add_title_slide(title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle
    return slide


def add_bullets(title, bullets, note=None):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if isinstance(b, tuple):
            text, level = b
        else:
            text, level = b, 0
        p.text = text
        p.level = level
    if note:
        tx = slide.shapes.add_textbox(Inches(0.5), Inches(6.6), Inches(12.5), Inches(0.6))
        p = tx.text_frame.paragraphs[0]
        p.text = f"話すポイント: {note}"
        p.font.size = Pt(11)
        p.font.color.rgb = GRAY


def add_loop_diagram():
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "3. 仕事のループとは何か"
    labels = ["情報を受け取る", "処理する", "保存・反映する", "次の処理へつなぐ"]
    positions = [(1.0,2.0),(4.0,2.0),(7.2,2.0),(10.0,2.0)]
    for i,(x,y) in enumerate(positions):
        shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(2.3), Inches(1.0))
        shp.fill.solid(); shp.fill.fore_color.rgb = BG
        shp.line.color.rgb = ACCENT
        shp.text_frame.text = labels[i]
        shp.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        if i < len(positions)-1:
            ar = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x+2.35), Inches(2.25), Inches(0.5), Inches(0.5))
            ar.fill.solid(); ar.fill.fore_color.rgb = ACCENT
            ar.line.color.rgb = ACCENT
    note = slide.shapes.add_textbox(Inches(1.0), Inches(3.8), Inches(11.0), Inches(1.2))
    note.text_frame.text = "従来AIは③④で人間の手作業（コピペ・保存・反映）が必要で、ループが閉じなかった。"


def add_before_after():
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "6. Before / After（体験ベース）"
    left = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.6), Inches(5.7), Inches(4.6))
    right = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.8), Inches(1.6), Inches(5.7), Inches(4.6))
    for shp, title in [(left, "Before（チャットAI中心）"),(right,"After（エージェントAI）")]:
        shp.fill.solid(); shp.fill.fore_color.rgb = BG
        shp.line.color.rgb = NAVY
        tf = shp.text_frame
        tf.clear()
        p = tf.paragraphs[0]; p.text = title; p.font.bold = True; p.font.size = Pt(18); p.font.color.rgb = NAVY
    lb = left.text_frame
    for t in ["・回答を受け取る", "・人間がコピペ", "・保存先を決める", "・別アプリへ反映", "・またAIに戻る"]:
        p = lb.add_paragraph(); p.text = t; p.level = 0
    rb = right.text_frame
    for t in ["・AIがファイルを直接操作", "・PC/アプリを連続操作", "・Webで取得→整理→反映", "・人間は確認と承認中心", "・ループがほぼ閉じる"]:
        p = rb.add_paragraph(); p.text = t; p.level = 0


def add_three_capabilities():
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "5. 何が変わったか（3つの能力）"
    caps = [
        ("ストレージアクセス", "読む / 書く / 更新 / 削除"),
        ("PC・アプリ操作", "ツール実行 / アプリ連携"),
        ("Webアクセス・操作", "取得 / 抽出 / 入力 / 操作"),
    ]
    for i,(t,s) in enumerate(caps):
        x = 0.9 + i*4.1
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.0), Inches(3.6), Inches(2.2))
        box.fill.solid(); box.fill.fore_color.rgb = BG
        box.line.color.rgb = ACCENT
        tf = box.text_frame
        tf.clear()
        p = tf.paragraphs[0]; p.text=t; p.font.bold=True; p.font.size=Pt(16); p.font.color.rgb = NAVY
        p = tf.add_paragraph(); p.text=s; p.font.size=Pt(13)


add_title_slide(
    "非エンジニア目線で考える\nAIエージェントの本質",
    "有志参加AI勉強会（8分 / 10枚）\n2026-03-22"
)

add_bullets(
    "2. 先に結論",
    [
        "AIエージェントの本質は『よく答えること』ではない",
        "本質は『仕事のループを閉じること』",
        "答えるAI → 完遂するAI への構造変化が起きている",
    ],
    note="最初にゴールを示して、以降のスライドをこの軸で読む。",
)

add_loop_diagram()

add_bullets(
    "4. 従来AIの制約（なぜ閉じなかったか）",
    [
        "制約A: インターフェース外に出られない",
        ("→ 別アプリ反映は人間のコピペ前提",1),
        "制約B: ストレージ操作権限がない",
        ("→ 作成/更新/削除をAI単独で実行できない",1),
    ],
)

add_three_capabilities()

add_before_after()

add_bullets(
    "7. 実体験: しずく + OpenClaw",
    [
        "ローカルPCでファイル作成・編集・保存を連続実行",
        "ブラウザ操作で検索→抽出→報告まで自動化",
        "定期収集（為替/指標）をcron・heartbeatで自動運用",
        "人間は『指示・承認・最終判断』へ役割がシフト",
    ],
)

add_bullets(
    "8. 非エンジニアにとっての課題",
    [
        "設定実装の難しさ（使える状態へ落とし込む）",
        "セキュリティ設計（権限とリスクのトレードオフ）",
        "便利にするほど攻撃面・誤操作リスクも増える",
    ],
)

add_bullets(
    "9. 実践アプローチ（やさしい導入順）",
    [
        "段階1: 読み取り中心で開始（低リスク）",
        "段階2: 限定書き込みを許可（対象を明確化）",
        "段階3: 自動実行を導入（cron/heartbeat）",
        "常時: バックアップ・ログ・停止手段を先に用意",
    ],
)

add_bullets(
    "10. まとめ",
    [
        "画期性は『IQ向上』だけでなく『I/Oと権限の進化』",
        "3能力（ストレージ・PC/アプリ・Web）でループが閉じる",
        "次の焦点は『非エンジニアでも安全に使える設計』",
        "一言: AIエージェントは“答える”から“やり切る”へ",
    ],
)

prs.save(OUT)
print(OUT)
